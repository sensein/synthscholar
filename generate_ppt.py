"""PRISMA Review Agent — Architecture PPT (15 slides)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
BLUE   = RGBColor(0x03, 0x69, 0xA1)
LBLUE  = RGBColor(0x7D, 0xD3, 0xFC)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
LGREEN = RGBColor(0x86, 0xEF, 0xAC)
YELLOW = RGBColor(0xCA, 0x8A, 0x04)
LYELL  = RGBColor(0xFD, 0xE6, 0x8A)
ORANGE = RGBColor(0xC2, 0x41, 0x0C)
LORNG  = RGBColor(0xFE, 0xD7, 0xAA)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
LPURP  = RGBColor(0xDD, 0xD6, 0xFE)
RED    = RGBColor(0xBE, 0x12, 0x3C)
TEAL   = RGBColor(0x0D, 0x94, 0x88)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xCC, 0xD6, 0xE0)
DGRAY  = RGBColor(0x15, 0x2A, 0x3C)
SLATE  = RGBColor(0x47, 0x55, 0x69)

DIAG = "/Users/tekrajchhetri/Documents/research_codes_papers_writing/synthscholar/PRISMA_Agent_Architecture.png"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width; H = prs.slide_height

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(s, color=NAVY):
    sh = s.shapes.add_shape(1,0,0,W,H)
    sh.fill.solid(); sh.fill.fore_color.rgb=color
    sh.line.fill.background()

def rect(s,x,y,w,h,fill,border=None,bw=Pt(1.5)):
    sh=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if border: sh.line.color.rgb=border; sh.line.width=bw
    else: sh.line.fill.background()
    return sh

def txt(s,text,x,y,w,h,size=12,bold=False,color=WHITE,
        align=PP_ALIGN.LEFT,italic=False):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold
    r.font.color.rgb=color; r.font.italic=italic

def hdr(s,title,sub=None):
    rect(s,0,0,13.33,0.65,BLUE)
    txt(s,title,0.3,0.08,12.7,0.5,size=22,bold=True,color=NAVY,align=PP_ALIGN.CENTER)
    if sub:
        txt(s,sub,0.3,0.72,12.7,0.38,size=11,color=LBLUE,align=PP_ALIGN.CENTER)

def bar(s,color=BLUE):
    sh=s.shapes.add_shape(1,0,H-Inches(0.07),W,Inches(0.07))
    sh.fill.solid(); sh.fill.fore_color.rgb=color
    sh.line.fill.background()

def num(s,n):
    txt(s,str(n),12.85,7.15,0.35,0.28,size=9,color=LBLUE,align=PP_ALIGN.RIGHT)

def card(s,x,y,w,h,title,bullets,tc=BLUE,bs=11):
    rect(s,x,y,w,h,DGRAY,tc,Pt(1.2))
    txt(s,title,x+0.1,y+0.06,w-0.2,0.36,size=13,bold=True,color=tc)
    if bullets:
        body="\n".join(f"  {b}" for b in bullets)
        txt(s,body,x+0.1,y+0.46,w-0.2,h-0.54,size=bs,color=LGRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s)
rect(s,0,0,13.33,0.72,BLUE)
txt(s,"PRISMA Review Agent",0.3,0.1,12.7,0.55,size=26,bold=True,color=NAVY,align=PP_ALIGN.CENTER)
txt(s,"Automated Systematic Literature Review  —  PRISMA 2020",
    0.3,1.0,12.7,0.7,size=36,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"System Architecture  |  9 LLM Agents  |  15-Step Pipeline  |  PRISMA Flow",
    0.3,1.85,12.7,0.45,size=17,color=LBLUE,align=PP_ALIGN.CENTER)
rect(s,4.5,2.6,4.3,0.05,BLUE)
txt(s,"Tek Raj Chhetri   |   v1.0.0   |   Apache 2.0",
    0.3,2.75,12.7,0.4,size=15,color=LGRAY,align=PP_ALIGN.CENTER)

pills=[("ReviewProtocol Input",BLUE),("9 Pydantic-AI Agents",GREEN),
       ("PubMed + bioRxiv",YELLOW),("PRISMA 2020 Output",ORANGE),
       ("GRADE + RoB Assessment",PURPLE)]
for i,(label,clr) in enumerate(pills):
    x=0.4+i*2.5
    rect(s,x,3.5,2.35,0.48,DGRAY,clr,Pt(1.5))
    txt(s,label,x+0.08,3.55,2.19,0.38,size=11,bold=True,color=clr,align=PP_ALIGN.CENTER)

txt(s,"ReviewProtocol (PICO + criteria)  ->  Search Agent  ->  PubMed/bioRxiv Tools  ->  Two-Stage Screening  ->  Evidence/RoB/Data  ->  Parallel Synthesis  ->  PRISMA Report",
    0.3,4.32,12.7,0.52,size=13,color=LGRAY,align=PP_ALIGN.CENTER)
txt(s,"Domain-agnostic  |  11 RoB tools  |  GRADE certainty  |  Any OpenRouter model  |  SQLite cache",
    0.3,5.05,12.7,0.42,size=13,color=LBLUE,align=PP_ALIGN.CENTER)
bar(s); num(s,1)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — Architecture Diagram
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s,RGBColor(0xF0,0xF4,0xF8)); num(s,2)
rect(s,0,0,13.33,0.58,NAVY)
txt(s,"SYSTEM ARCHITECTURE DIAGRAM",0.3,0.08,12.7,0.44,size=20,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
if os.path.exists(DIAG):
    s.shapes.add_picture(DIAG,Inches(0.08),Inches(0.65),Inches(13.18),Inches(6.75))
else:
    txt(s,"Run make_diagram.py first",2,3,9,1,size=16,color=RED,align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — ReviewProtocol & Purpose
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,3)
hdr(s,"REVIEW PROTOCOL — The System's Input",
    "All domain knowledge injected at runtime — system is fully domain-agnostic")

rect(s,0.2,1.2,6.1,5.8,DGRAY,PURPLE,Pt(2))
txt(s,"ReviewProtocol  (models.py)",0.35,1.28,5.8,0.38,size=14,bold=True,color=PURPLE)
fields=[
    ("title","Research question / review title"),
    ("objective","Detailed objective statement"),
    ("pico_population","Who — target population"),
    ("pico_intervention","What — intervention being studied"),
    ("pico_comparison","Comparator — control or alternative"),
    ("pico_outcome","What is measured — primary outcome"),
    ("inclusion_criteria","Full inclusion rules (free text)"),
    ("exclusion_criteria","Full exclusion rules (free text)"),
    ("databases","PubMed, bioRxiv (extensible)"),
    ("rob_tool","1 of 11 RoB tools: RoB 2, ROBINS-I, NOS..."),
    ("date_range_start/end","Optional date restriction"),
    ("registration","PROSPERO or other registry number"),
]
for i,(field,desc) in enumerate(fields):
    y=1.78+i*0.4
    rect(s,0.3,y,5.9,0.36,NAVY,PURPLE,Pt(0.6))
    txt(s,field,0.42,y+0.04,2.2,0.28,size=10,bold=True,color=PURPLE)
    txt(s,desc,2.65,y+0.04,3.5,0.28,size=10,color=LGRAY)

rect(s,6.5,1.2,6.6,5.8,DGRAY,BLUE,Pt(2))
txt(s,"Why domain-agnostic design?",6.65,1.28,6.3,0.38,size=14,bold=True,color=BLUE)
reasons=[
    "No hardcoded medical knowledge anywhere",
    "Same codebase conducts RCT reviews in cardiology,",
    "  cohort studies in neuroscience, qualitative studies",
    "  in education — without any code changes",
    "",
    "How it works:",
    "Every LLM agent receives protocol fields via",
    "  dynamic @agent.system_prompt context at call time",
    "Agent sees: question, PICO, criteria, RoB tool",
    "Agent does NOT hardcode any disease/domain rules",
    "",
    "Example: rob_agent uses protocol.rob_tool to select",
    "  domain list from ROB_DOMAINS dict:",
    "  RoB 2 -> 5 domains",
    "  ROBINS-I -> 7 domains",
    "  Newcastle-Ottawa -> 3 groups",
    "  QUADAS-2 -> 4 domains",
    "  ...11 tools total",
    "",
    "Outcome: one agent, eleven assessment frameworks,",
    "  zero code changes between reviews",
]
for i,r in enumerate(reasons):
    txt(s,r,6.65,1.75+i*0.28,6.3,0.26,size=9,color=LGRAY if not r.startswith("Why") and not r.startswith("How") and not r.startswith("Outcome") else BLUE)
bar(s)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — Two-Stage Screening Design
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,4)
hdr(s,"TWO-STAGE SCREENING — Same Agent, Opposite Bias Policies",
    "PRISMA best practice: maximise recall at abstract stage, maximise precision at full-text stage")

rect(s,0.2,1.2,6.1,5.8,DGRAY,BLUE,Pt(2))
txt(s,"Stage A: Title / Abstract Screening",0.35,1.28,5.8,0.38,size=13,bold=True,color=BLUE)
txt(s,"Agent 2: screening_agent  |  Batch size: 15",0.35,1.7,5.8,0.28,size=10,color=LBLUE,italic=True)
rect(s,0.3,2.08,5.9,0.42,RGBColor(0x03,0x69,0xA1),BLUE,Pt(0))
txt(s,"POLICY: INCLUSIVE — when in doubt, INCLUDE",0.42,2.12,5.7,0.34,size=12,bold=True,color=WHITE)
ta_lines=[
    "Input: title + abstract of each article",
    "Eligibility criteria injected from protocol",
    "Decision per article: INCLUDE or EXCLUDE",
    "  + reason string + relevance_score 0-1",
    "On agent error: auto-include entire batch",
    "  (never lose an article to a transient failure)",
    "",
    "PRISMA rationale:",
    "  Excluding a relevant article at abstract stage",
    "  is more harmful than including an irrelevant one.",
    "  Full-text screening is the precision filter.",
    "  Abstract screening is the recall filter.",
    "",
    "Flow tracking: flow.excluded_title_abstract",
    "  = count of excluded articles here",
]
for i,l in enumerate(ta_lines):
    txt(s,l,0.35,2.6+i*0.3,5.8,0.28,size=9.5,color=LGRAY if not l.startswith("PRISMA") else LBLUE)

rect(s,6.5,1.2,6.6,5.8,DGRAY,ORANGE,Pt(2))
txt(s,"Stage B: Full-text Eligibility Screening",6.65,1.28,6.3,0.38,size=13,bold=True,color=ORANGE)
txt(s,"Agent 2: screening_agent  |  Batch size: 10",6.65,1.7,6.3,0.28,size=10,color=LORNG,italic=True)
rect(s,6.58,2.08,6.34,0.42,RGBColor(0xC2,0x41,0x0C),ORANGE,Pt(0))
txt(s,"POLICY: STRICT — apply criteria rigorously",6.7,2.12,6.1,0.34,size=12,bold=True,color=WHITE)
ft_lines=[
    "Input: full text (up to 12,000 chars from PMC)",
    "  + title + abstract",
    "Articles without PMC ID: auto-include",
    "  (cannot assess eligibility without text)",
    "Decision: INCLUDE or EXCLUDE",
    "  + specific exclusion reason",
    "Batch size reduced to 10 (full text = more tokens)",
    "",
    "Exclusion reasons accumulated into:",
    "  flow.excluded_reasons  (dict: reason -> count)",
    "  Top 8 reasons shown in PRISMA flow diagram",
    "",
    "PRISMA rationale:",
    "  Final gate before synthesis.",
    "  Only studies clearly satisfying ALL inclusion",
    "  criteria and violating NO exclusion criteria",
    "  proceed to evidence extraction and synthesis.",
]
for i,l in enumerate(ft_lines):
    txt(s,l,6.65,2.6+i*0.3,6.3,0.28,size=9.5,color=LGRAY if not l.startswith("PRISMA") else LORNG)
bar(s)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — The 9 Agents
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,5)
hdr(s,"THE 9 LLM AGENTS  (pydantic-ai + OpenRouter)",
    "Each agent has a single PRISMA responsibility  |  defer_model_check=True -> any model at runtime")

agents=[
    ("Agent 1","search_strategy_agent",YELLOW,
     "In: ReviewProtocol\nOut: SearchStrategy",
     ["2-5 PubMed queries (MeSH, Boolean, field tags)",
      "2-3 bioRxiv plain-text queries",
      "mesh_terms[] for methods documentation",
      "rationale for search strategy"]),
    ("Agent 2","screening_agent",BLUE,
     "In: Article batch + criteria\nOut: ScreeningBatchResult",
     ["Two-stage: T/A (batch 15, inclusive) + FT (batch 10, strict)",
      "Per article: decision, reason, relevance_score",
      "Auto-include entire batch on error",
      "Same agent, different bias from protocol stage"]),
    ("Agent 3","rob_agent",ORANGE,
     "In: Article + rob_tool\nOut: RiskOfBiasResult",
     ["Per-article, per-domain judgment: LOW/SOME/HIGH/UNCLEAR",
      "11 tools: RoB 2, ROBINS-I, NOS, QUADAS-2, CASP, JBI...",
      "Overall judgment + justification per domain",
      "Tool + domains injected from ReviewProtocol.rob_tool"]),
    ("Agent 4","data_extraction_agent",GREEN,
     "In: Article + data_items\nOut: StudyDataExtraction",
     ["Optional (--extract-data flag)",
      "study_design, sample_size, population",
      "intervention, outcomes, key_findings",
      "effect_measures (ORs, MDs, HRs if reported)"]),
    ("Agent 5","synthesis_agent",PURPLE,
     "In: articles (top 25) + evidence (top 20)\nOut: str Markdown narrative",
     ["Thematic organisation, not study-by-study",
      "Cite: (Author et al., Year; PMID: XXXXX)",
      "Explicitly note contradictions between studies",
      "PRISMA 2020 Results section format"]),
    ("Agent 6","grade_agent",TEAL,
     "In: outcome name + article list\nOut: GRADEAssessment",
     ["Per-outcome certainty rating: HIGH/MODERATE/LOW/VERY LOW",
      "5 domains: RoB, inconsistency, indirectness,",
      "  imprecision, publication bias",
      "Run x3 outcomes in parallel via asyncio.gather()"]),
    ("Agent 7","bias_summary_agent",YELLOW,
     "In: included article list\nOut: str (bias summary)",
     ["Overall quality across all studies",
      "Common methodological weaknesses",
      "Publication bias concerns",
      "Confidence in body of evidence"]),
    ("Agent 8","limitations_agent",RED,
     "In: flow text + article list\nOut: str (limitations)",
     ["Search scope and database limitations",
      "Full-text retrieval gaps (paywall)",
      "AI-assisted screening caveat (required)",
      "Heterogeneity and meta-analysis limitations"]),
    ("Agent 9","evidence_extraction_agent",GREEN,
     "In: articles batch (5)\nOut: BatchEvidenceExtraction",
     ["2-5 EvidenceSpan per article",
      "text, claim, section, relevance_score, is_quantitative",
      "Dedup: Jaccard word overlap >= 0.70 -> remove",
      "Cap: top 30 spans retained for synthesis"]),
]

for i,(badge,name,color,io,bullets) in enumerate(agents):
    col=i%3; row=i//3
    x=0.2+col*4.35; y=1.2+row*2.05
    bh=1.92
    rect(s,x,y,4.2,bh,DGRAY,color,Pt(1.2))
    rect(s,x,y,0.88,0.4,color)
    txt(s,badge,x+0.02,y+0.04,0.84,0.32,size=10,bold=True,color=NAVY,align=PP_ALIGN.CENTER)
    txt(s,name,x+0.96,y+0.05,3.16,0.3,size=10,bold=True,color=color)
    txt(s,io,x+0.1,y+0.45,4.0,0.3,size=8.5,color=LBLUE,italic=True)
    body="\n".join(f"  {b}" for b in bullets)
    txt(s,body,x+0.1,y+0.78,4.0,bh-0.85,size=8.5,color=LGRAY)
bar(s)

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Pipeline Steps 1-8
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,6)
hdr(s,"15-STEP PIPELINE  (Steps 1–8: Protocol to Screened Articles)",
    "Sequential execution with SQLite caching at every data collection step")

steps=[
    ("01","Protocol Input","User / CLI",PURPLE,
     "ReviewProtocol: title, PICO, criteria, databases, rob_tool, date range. All domain knowledge lives here."),
    ("02","Search Strategy","Agent 1",YELLOW,
     "protocol -> search_strategy_agent -> SearchStrategy: 2-5 PubMed queries + 2-3 bioRxiv queries + mesh_terms"),
    ("03","PubMed Search","PubMed Tool (NCBI)",GREEN,
     "Each pubmed_query -> esearch.fcgi (PMID list) -> efetch.fcgi XML (batch 50) -> Article[]. Cached 72h."),
    ("04","bioRxiv Search","bioRxiv Tool",GREEN,
     "Each biorxiv_query -> REST JSON API -> word-overlap filter (>=2 words) -> Article[source=biorxiv]. 180-day window."),
    ("05","Related Articles","elink Tool",TEAL,
     "Top 8 PubMed PMIDs -> elink neighbor_score -> up to 15 related PMIDs per depth. Fetch and add. Depth 1..N."),
    ("06","Citation Hops","elink Tool",TEAL,
     "Top 5 PMIDs -> backward (neighbor_score) + forward (citedin) -> combine unique -> fetch new articles. Hops 1..max."),
    ("07","Deduplication","Pipeline logic",SLATE,
     "Key = DOI.lower() if doi else pmid. Keep first occurrence per key. flow.duplicates_removed = total - unique."),
    ("08","Title/Abstract Screen","Agent 2 (INCLUSIVE)",BLUE,
     "Batch 15, inclusive bias. Per article: INCLUDE/EXCLUDE + reason + relevance_score. Auto-include on error."),
]

for i,(num_s,title,source,color,why) in enumerate(steps):
    col=i%2; row=i//2
    x=0.2+col*6.55; y=1.2+row*1.5
    rect(s,x,y,6.38,1.4,DGRAY,color,Pt(1.1))
    txt(s,f"Step {num_s}  {title}",x+0.1,y+0.05,4.5,0.35,size=12,bold=True,color=color)
    txt(s,f"[{source}]",x+0.1,y+0.42,6.1,0.25,size=10,color=LBLUE,italic=True)
    txt(s,why,x+0.1,y+0.72,6.1,0.62,size=9.5,color=LGRAY)
bar(s)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Pipeline Steps 9-15
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,7)
hdr(s,"15-STEP PIPELINE  (Steps 9–15: Full-text to Final Report)",
    "Article-level analysis then parallel synthesis via asyncio.gather()")

steps2=[
    ("09","PMC Full-text Fetch","PMC Tool (NCBI XML)",GREEN,
     "Filter ta_included by pmc_id -> efetch?db=pmc XML -> extract <body> -> strip tags -> 12k chars max. 10/call."),
    ("10","Full-text Eligibility","Agent 2 (STRICT)",BLUE,
     "Batch 10, strict bias. Exclusion reasons -> flow.excluded_reasons{}. Auto-include if no full text."),
    ("11","Evidence Extraction","Agent 9 (batch 5)",ORANGE,
     "2-5 EvidenceSpan/article: text, claim, section, score. Jaccard dedup >= 0.70. Top 30 retained."),
    ("12","Data Extraction","Agent 4 (per-article, optional)",GREEN,
     "study_design, sample_size, population, intervention, outcomes, key_findings, effect_measures. --extract-data flag."),
    ("13","Risk of Bias","Agent 3 (per-article)",ORANGE,
     "Per-domain judgment (LOW/SOME/HIGH) using protocol.rob_tool. Overall judgment + justification. 11 tools supported."),
    ("14","Parallel Synthesis","asyncio.gather() — Agents 5,6,7,8",PURPLE,
     "Simultaneously: synthesis narrative (Agent 5) + bias summary (Agent 7) + GRADE per outcome (Agent 6) + limitations (Agent 8)."),
    ("15","Assemble Result","Pipeline",TEAL,
     "Construct PRISMAReviewResult: protocol, flow, included_articles, screening_log, evidence, synthesis, GRADE, timestamp."),
]

for i,(num_s,title,source,color,why) in enumerate(steps2):
    col=i%2; row=i//2
    x=0.2+col*6.55; y=1.2+row*1.73
    bh=1.6
    rect(s,x,y,6.38,bh,DGRAY,color,Pt(1.1))
    txt(s,f"Step {num_s}  {title}",x+0.1,y+0.05,4.5,0.35,size=12,bold=True,color=color)
    txt(s,f"[{source}]",x+0.1,y+0.44,6.1,0.25,size=10,color=LBLUE,italic=True)
    txt(s,why,x+0.1,y+0.74,6.1,0.78,size=9.5,color=LGRAY)
bar(s)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Parallel Synthesis Step
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,8)
hdr(s,"PARALLEL SYNTHESIS  (asyncio.gather — Step 14)",
    "All four synthesis tasks are independent — run simultaneously to save time")

rect(s,0.2,1.2,12.95,0.52,DGRAY,PURPLE,Pt(1.5))
txt(s,"asyncio.gather(run_synthesis(), run_bias_summary(), run_grade() x3, run_limitations(), return_exceptions=True)",
    0.35,1.28,12.65,0.38,size=11,bold=True,color=PURPLE)

par_agents=[
    ("Agent 5","synthesis_agent",PURPLE,
     "Input",
     "included_articles[:25]  +  evidence_spans[:20]  +  PRISMA flow summary text",
     "Output",
     "str  (Markdown narrative synthesis)",
     ["Thematic organisation — groups findings across studies",
      "Citation format: (Author et al., Year; PMID: XXXXX)",
      "Explicitly note contradictions between studies",
      "Becomes the Results / Synthesis section of the report",
      "NEVER fabricates PMIDs or statistics"]),
    ("Agent 7","bias_summary_agent",YELLOW,
     "Input","included article list with RoB assessments",
     "Output","str  (overall quality summary)",
     ["Proportion of studies at low/high/unclear risk",
      "Most common methodological weaknesses",
      "Publication bias concerns",
      "Confidence in the overall body of evidence"]),
    ("Agent 6","grade_agent  (x3 outcomes)",TEAL,
     "Input","outcome name + list of relevant studies",
     "Output","GRADEAssessment  (Pydantic model)",
     ["5 GRADE domains: risk of bias, inconsistency,",
      "  indirectness, imprecision, publication bias",
      "Per-domain: downgrade or no concerns",
      "Overall: HIGH / MODERATE / LOW / VERY LOW",
      "Appears as GRADE table in Methods section"]),
    ("Agent 8","limitations_agent",RED,
     "Input","PRISMA flow text + included article list",
     "Output","str  (limitations section)",
     ["Search scope (PubMed + bioRxiv only)",
      "Full-text retrieval gaps (paywall articles)",
      "AI-assisted screening caveat (required by PRISMA AI guidance)",
      "Language / date range restrictions",
      "Heterogeneity preventing meta-analysis"]),
]

for i,(badge,name,color,in_lbl,in_val,out_lbl,out_val,bullets) in enumerate(par_agents):
    col=i%2; row=i//2
    x=0.2+col*6.55; y=2.0+row*2.65
    bh=2.5
    rect(s,x,y,6.38,bh,DGRAY,color,Pt(1.5))
    rect(s,x,y,0.9,0.4,color)
    txt(s,badge,x+0.04,y+0.04,0.82,0.32,size=10,bold=True,color=NAVY,align=PP_ALIGN.CENTER)
    txt(s,name,x+0.98,y+0.05,5.32,0.32,size=11,bold=True,color=color)
    txt(s,f"In: {in_val}",x+0.1,y+0.48,6.1,0.26,size=9,color=LBLUE)
    txt(s,f"Out: {out_val}",x+0.1,y+0.76,6.1,0.26,size=9,color=YELLOW)
    body="\n".join(f"  {b}" for b in bullets)
    txt(s,body,x+0.1,y+1.06,6.1,bh-1.14,size=9,color=LGRAY)
bar(s)

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Risk of Bias & GRADE
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,9)
hdr(s,"RISK OF BIAS & GRADE CERTAINTY",
    "11 RoB tools supported  |  GRADE certainty for up to 3 outcomes  |  all injected from ReviewProtocol")

rect(s,0.2,1.2,6.1,5.8,DGRAY,ORANGE,Pt(2))
txt(s,"Risk of Bias Assessment  (Agent 3)",0.35,1.28,5.8,0.38,size=14,bold=True,color=ORANGE)
rob_tools=[
    ("RoB 2","Randomised trials: 5 domains"),
    ("ROBINS-I","Non-randomised studies: 7 domains"),
    ("Newcastle-Ottawa","Case-control/cohort: 3 groups"),
    ("QUADAS-2","Diagnostic accuracy: 4 domains"),
    ("CASP","Qualitative: 10 items"),
    ("JBI","Prevalence/cross-sectional: 7 items"),
    ("Murad","Case reports/series: 4 items"),
    ("Jadad","RCT quality: 5-point scale"),
    ("ROBINS-E","Environmental exposures"),
    ("SYRCLE","Animal studies"),
    ("MINORS","Non-randomised surgical"),
]
for i,(tool,desc) in enumerate(rob_tools):
    y=1.76+i*0.44
    rect(s,0.3,y,5.9,0.38,NAVY,ORANGE,Pt(0.6))
    txt(s,tool,0.42,y+0.04,2.0,0.3,size=10,bold=True,color=ORANGE)
    txt(s,desc,2.45,y+0.04,3.7,0.3,size=10,color=LGRAY)
txt(s,"Per domain: LOW / SOME CONCERNS / HIGH / UNCLEAR\nOverall judgment + justification per domain",
    0.35,6.58,5.8,0.32,size=9.5,color=LORNG)

rect(s,6.5,1.2,6.6,5.8,DGRAY,TEAL,Pt(2))
txt(s,"GRADE Certainty Assessment  (Agent 6)",6.65,1.28,6.3,0.38,size=14,bold=True,color=TEAL)
grade_items=[
    ("Risk of bias","RoB assessments across included studies"),
    ("Inconsistency","Variation in findings across studies"),
    ("Indirectness","How directly evidence addresses question"),
    ("Imprecision","Width of confidence intervals / sample size"),
    ("Publication bias","Selective reporting, funnel plot asymmetry"),
]
txt(s,"5 GRADE Domains:",6.65,1.76,6.3,0.3,size=11,color=WHITE,bold=True)
for i,(domain,desc) in enumerate(grade_items):
    y=2.12+i*0.58
    rect(s,6.58,y,6.34,0.52,NAVY,TEAL,Pt(0.8))
    txt(s,domain,6.7,y+0.04,2.5,0.24,size=10,bold=True,color=TEAL)
    txt(s,desc,9.22,y+0.04,3.6,0.24,size=9.5,color=LGRAY)

txt(s,"Overall Certainty:",6.65,5.1,6.3,0.3,size=11,bold=True,color=WHITE)
for i,(cert,desc,clr) in enumerate([
    ("HIGH","Strong confidence in estimate",GREEN),
    ("MODERATE","Moderate confidence; likely close",TEAL),
    ("LOW","Limited confidence; substantially different",YELLOW),
    ("VERY LOW","Very little confidence",RED),
]):
    y=5.45+i*0.43
    rect(s,6.58,y,1.6,0.36,clr)
    txt(s,cert,6.62,y+0.04,1.52,0.28,size=9,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s,desc,8.22,y+0.04,4.8,0.28,size=9,color=LGRAY)
bar(s)

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — Data Collection Tools
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,10)
hdr(s,"DATA COLLECTION TOOLS  (clients.py)",
    "HTTP clients called by pipeline  |  all results cached in SQLite (72h TTL)  |  no pydantic-ai @agent.tool")

tools=[
    ("PubMed Search Tool",BLUE,[
        "API: NCBI E-utilities REST",
        "esearch.fcgi?db=pubmed -> PMID list (JSON)",
        "efetch.fcgi?db=pubmed&retmode=xml -> PubMed XML",
        "Batch: 50 PMIDs per fetch call",
        "Parse: regex on <ArticleTitle> <AbstractText> <Author>",
        "  <ISOAbbreviation> <PubDate> <ArticleId> <MeshHeading>",
        "Rate: 0.35s between calls (3 req/s; 10 with NCBI key)",
        "Cache: search namespace (PMID lists) + article namespace"]),
    ("PMC Full-text Tool",GREEN,[
        "API: NCBI E-utilities (same client)",
        "efetch.fcgi?db=pmc&id=<pmc_id>&rettype=xml",
        "Extract <body> XML tag content",
        "Strip all HTML/XML tags (regex)",
        "Normalise whitespace, truncate to 12,000 chars",
        "Max 10 per call; only open-access articles",
        "Cache: fulltext namespace",
        "NO PDF — direct structured XML only"]),
    ("bioRxiv Tool",YELLOW,[
        "API: bioRxiv REST API (JSON, not XML)",
        "https://api.biorxiv.org/details/biorxiv/{start}/{end}/{cursor}/30",
        "Default window: last 180 days (--biorxiv-days)",
        "Pagination: 30 per page, up to 4 pages (120 candidates)",
        "Filter: word-overlap >= 2 non-trivial words",
        "Article.pmid = 'biorxiv_{doi_suffix}'",
        "Article.source = 'biorxiv'",
        "Cache: biorxiv namespace"]),
    ("Related Articles Tool",TEAL,[
        "API: NCBI elink.fcgi",
        "LinkName=pubmed_pubmed_related (neighbor_score)",
        "Seeds: top 8 PubMed PMIDs per depth level",
        "Max 15 related per call",
        "Depth: 1..related_depth (default: 1)",
        "Article.source = 'related_N' (depth number)",
        "Cache: related namespace"]),
    ("Citation Hops Tool",ORANGE,[
        "Backward: elink neighbor_score (same as Related)",
        "Forward: elink?LinkName=pubmed_pubmed_citedin",
        "Seeds: top 5 PubMed PMIDs",
        "Combine backward + forward unique PMIDs",
        "Article.hop_level + Article.parent_id set",
        "Article.source = 'hop_N'",
        "Max hops: --hops parameter (default: 1, max: 4)"]),
]

for i,(name,color,lines) in enumerate(tools):
    col=i%2 if i<4 else 0
    row=i//2
    if i==4:
        x=0.2; y=1.2+row*2.95; bw=6.38
    else:
        x=0.2+col*6.55; y=1.2+row*2.95; bw=6.38
    bh=2.75 if i<4 else 2.75
    rect(s,x,y,bw,bh,DGRAY,color,Pt(1.3))
    txt(s,name,x+0.1,y+0.06,bw-0.2,0.36,size=13,bold=True,color=color)
    body="\n".join(f"  {l}" for l in lines)
    txt(s,body,x+0.1,y+0.5,bw-0.2,bh-0.58,size=9.5,color=LGRAY)

if len(tools)>4:
    rect(s,6.8,1.2+2*2.95,6.35,2.75,DGRAY,PURPLE,Pt(1.3))
    txt(s,"SQLite Cache  (cache.py)",6.95,1.26+5.9,6.1,0.36,size=13,bold=True,color=PURPLE)
    clines=["TTL: 72 hours  |  5 namespaces",
            "search: PubMed PMID lists","article: individual Article records",
            "related: elink PMID lists","fulltext: PMC body text",
            "biorxiv: bioRxiv results","Key: SHA256(namespace:identifier)",
            "Lazy expiry: check on read, delete if expired"]
    body="\n".join(f"  {l}" for l in clines)
    txt(s,body,6.95,1.68+5.9,6.1,2.1,size=9.5,color=LGRAY)
bar(s)

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — Data Models
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,11)
hdr(s,"DATA MODELS  (models.py — 15+ Pydantic v2 Classes)",
    "Runtime-validated typed schemas for every input, intermediate state, and output")

groups=[
    ("INPUT",YELLOW,["ReviewProtocol",
        "  title, objective","  pico_population/intervention/comparison/outcome",
        "  inclusion_criteria, exclusion_criteria",
        "  databases, rob_tool, date_range","  registration"]),
    ("SEARCH",BLUE,["SearchStrategy",
        "  pubmed_queries: list[str]","  biorxiv_queries: list[str]",
        "  mesh_terms: list[str]","  rationale: str"]),
    ("ARTICLE",GREEN,["Article",
        "  pmid, title, abstract","  authors, journal, year",
        "  doi, pmc_id, mesh_terms","  full_text (after PMC fetch)",
        "  extracted_data: StudyDataExtraction","  risk_of_bias: RiskOfBiasResult",
        "  inclusion_status: InclusionStatus","  source, hop_level, parent_id"]),
    ("SCREENING",ORANGE,["ScreeningBatchResult",
        "  decisions: list[ScreeningDecision]",
        "ScreeningDecision:",
        "  index, decision (INCLUDE/EXCLUDE)","  reason: str",
        "  relevance_score: float",
        "InclusionStatus (Enum):",
        "  PENDING / INCLUDED / EXCLUDED"]),
    ("EVIDENCE",TEAL,["EvidenceSpan",
        "  text, paper_pmid, claim","  section, relevance_score",
        "BatchEvidenceExtraction",
        "  articles: list[ArticleEvidenceExtraction]",
        "ExtractedEvidenceItem:",
        "  quote, claim, section","  relevance (0-1), is_quantitative"]),
    ("RoB + GRADE",PURPLE,["RiskOfBiasResult",
        "  assessments: list[RoBDomainAssessment]","  overall: RoBJudgment",
        "  summary: str",
        "GRADEAssessment",
        "  outcome, domains{}","  overall_certainty: GRADECertainty",
        "  summary: str"]),
    ("DATA EXTRACT",RED,["StudyDataExtraction",
        "  study_design, sample_size","  population, intervention",
        "  outcomes, key_findings","  effect_measures (optional)"]),
    ("RESULT",BLUE,["PRISMAReviewResult",
        "  protocol, flow, included_articles","  screening_log, evidence_spans",
        "  synthesis_text, bias_assessment","  grade_assessments, limitations",
        "  timestamp",
        "PRISMAFlowCounts:",
        "  all PRISMA flow diagram numbers"]),
]

for i,(group,color,lines) in enumerate(groups):
    col=i%4; row=i//4
    x=0.2+col*3.28; y=1.2+row*2.95
    rect(s,x,y,3.15,2.78,DGRAY,color,Pt(1.3))
    txt(s,group,x+0.1,y+0.06,2.9,0.32,size=12,bold=True,color=color)
    body="\n".join(lines[:9])
    txt(s,body,x+0.1,y+0.44,2.9,2.27,size=8.5,color=LGRAY)
bar(s)

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — PRISMA Flow Counting
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,12)
hdr(s,"PRISMA 2020 FLOW COUNTING  (PRISMAFlowCounts)",
    "Every article transition tracked throughout pipeline — renders as PRISMA flow diagram in output report")

rect(s,0.2,1.2,12.95,0.46,DGRAY,BLUE,Pt(1.2))
txt(s,"Flow counters updated at each pipeline step  ->  rendered as PRISMA 2020 flow table in Markdown output",
    0.35,1.28,12.65,0.34,size=11,bold=True,color=BLUE)

flow_steps=[
    ("Step 3","db_pubmed","Articles from PubMed searches",GREEN),
    ("Step 4","db_biorxiv","Articles from bioRxiv searches",GREEN),
    ("Step 5","db_related","Articles from related expansion",TEAL),
    ("Step 6","db_hops","Articles from citation hops",TEAL),
    ("Steps 3-6","total_identified","Sum of all above (before dedup)",BLUE),
    ("Step 7","duplicates_removed","Removed by DOI/PMID deduplication",ORANGE),
    ("Step 7","after_dedup","Articles entering screening",BLUE),
    ("Step 8","screened_title_abstract","= after_dedup (all screened)",BLUE),
    ("Step 8","excluded_title_abstract","Excluded at T/A stage",RED),
    ("Steps 8-9","sought_fulltext","Passed T/A -> sought full text",BLUE),
    ("Step 9","not_retrieved","No PMC ID or retrieval failed",ORANGE),
    ("Step 10","assessed_eligibility","= sought_fulltext",BLUE),
    ("Step 10","excluded_eligibility","Excluded at full-text stage",RED),
    ("Step 10","excluded_reasons","Dict: reason -> count (top 8)",ORANGE),
    ("Step 10","included_synthesis","Final included articles",GREEN),
]

for i,(step,field,desc,color) in enumerate(flow_steps):
    col=i%2; row=i//2
    x=0.2+col*6.55; y=1.88+row*0.7
    rect(s,x,y,6.38,0.62,DGRAY,color,Pt(0.8))
    txt(s,step,x+0.1,y+0.06,1.0,0.22,size=9,bold=True,color=LGRAY)
    txt(s,field,x+1.15,y+0.06,2.5,0.22,size=10,bold=True,color=color)
    txt(s,desc,x+3.7,y+0.06,2.6,0.22,size=9.5,color=LGRAY)
bar(s)

# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — Export Formats & CLI
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,13)
hdr(s,"EXPORT FORMATS & CLI  (export.py + main.py)",
    "Three PRISMA-structured formats + full CLI control over every pipeline parameter")

for i,(fmt,color,use_case,sections) in enumerate([
    ("Markdown",BLUE,"PRISMA 2020 report, journal submission",
     ["Abstract (structured summary)","Introduction (rationale + PICO)",
      "Methods (criteria, queries, databases, RoB tool)","Results 3.1: PRISMA Flow table",
      "Results 3.2: Study Characteristics table","Results 3.3: Narrative Synthesis",
      "Results 3.4: Risk of Bias Summary","Results 3.5: GRADE Certainty table",
      "Discussion: Limitations","References (numbered with DOIs)","Appendix: evidence spans"]),
    ("JSON",YELLOW,"Programmatic use, REST APIs, databases",
     ["Full PRISMAReviewResult via model_dump_json()",
      "All nested Pydantic objects included","Protocol, flow, articles, screening_log",
      "evidence_spans, synthesis_text","bias_assessment, grade_assessments","limitations, timestamp"]),
    ("BibTeX",GREEN,"Zotero, Mendeley, LaTeX reference managers",
     ["@article entries for each included study",
      "Fields: title, author, journal, year, doi, pmid",
      "Files saved to prisma_results/{slug}_{timestamp}.bib"]),
]):
    x=0.2+i*4.35
    rect(s,x,1.2,4.2,3.5,DGRAY,color,Pt(1.5))
    txt(s,fmt,x+0.1,1.27,4.0,0.38,size=15,bold=True,color=color)
    txt(s,f"Use: {use_case}",x+0.1,1.7,4.0,0.26,size=10,color=LBLUE)
    body="\n".join(f"  {sec}" for sec in sections[:9])
    txt(s,body,x+0.1,2.0,4.0,2.55,size=9.5,color=LGRAY)

rect(s,0.2,4.9,12.95,2.42,DGRAY,BLUE,Pt(1.5))
txt(s,"CLI  (prisma-review entry point)",0.35,4.98,12.65,0.35,size=13,bold=True,color=BLUE)
cli=('prisma-review \\\n'
     '  --title "Effects of BDNF on hippocampal plasticity in major depression" \\\n'
     '  --population "Adults with MDD" --intervention "BDNF-targeted therapy" \\\n'
     '  --comparison "Placebo or standard care" --outcome "Depressive symptom reduction" \\\n'
     '  --inclusion "RCTs and cohort studies, adults, English" \\\n'
     '  --exclusion "Animal studies, review articles, case reports" \\\n'
     '  --rob-tool "RoB 2"  --model anthropic/claude-sonnet-4 \\\n'
     '  --max-results 25  --related-depth 2  --hops 2  --biorxiv-days 365 \\\n'
     '  --extract-data  --export md json bib  # or --interactive for guided setup')
txt(s,cli,0.35,5.34,12.65,1.92,size=10,color=GREEN)
bar(s)

# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — Design Decisions
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,14)
hdr(s,"KEY DESIGN DECISIONS & RATIONALE",
    "Every architectural choice reflects PRISMA methodology or engineering constraint")

decisions=[
    ("Agent-per-PRISMA-task\n(not a mega-agent)",BLUE,
     "Each PRISMA step gets its own agent: independent prompt tuning, isolated retry logic, typed output. "
     "Changing the screening bias policy only touches the screening agent. Adding GRADE for a new outcome adds one agent call."),
    ("Two-stage screening,\nopposite bias policies",GREEN,
     "PRISMA methodology requires recall-first, precision-second screening. Inclusive T/A ensures no relevant study is missed. "
     "Strict FT ensures only eligible studies enter synthesis. Same agent, different instruction from pipeline stage parameter."),
    ("defer_model_check=True\non all agents",YELLOW,
     "Model injected at runtime via build_model(). Researcher can use claude-haiku for screening (fast/cheap), "
     "claude-sonnet for synthesis (quality). Same agents work with any of 100+ OpenRouter models."),
    ("asyncio.gather() for\nStep 14 only",TEAL,
     "Synthesis, bias summary, GRADE, and limitations are independent tasks with no shared mutable state. "
     "Parallelising them with asyncio.gather saves time proportional to the slowest task instead of their sum."),
    ("SQLite cache, 72h TTL",ORANGE,
     "Researcher re-runs with different screening criteria or model. All 500 article fetches come from cache in milliseconds. "
     "72h (vs 48h in BioSynth) because PRISMA searches are longer sessions and article metadata changes infrequently."),
    ("Domain-agnostic by design",PURPLE,
     "All domain knowledge in ReviewProtocol. No hardcoded medical terminology, RoB domain lists, or outcome definitions. "
     "Injected as @agent.system_prompt context. Same binary, any medical/social/environmental domain."),
    ("Auto-include on\nagent failure",RED,
     "Pipeline must always produce a result. If screening batch fails, auto-include entire batch. "
     "Losing an article to a transient API error is a validity threat; over-inclusion is recoverable by human review."),
    ("Evidence dedup:\nJaccard >= 0.70",TEAL,
     "Same threshold as BioSynth. Removes paraphrased near-duplicates from the 30-span evidence set. "
     "Prevents the synthesis agent from seeing the same finding quoted five times and over-weighting it."),
]

for i,(title,color,text) in enumerate(decisions):
    col=i%2; row=i//4
    x=0.2+col*6.55; y=1.2+row*2.1+(i//2%2)*1.05
    rect(s,x,y,6.38,0.98,DGRAY,color,Pt(1.2))
    txt(s,title,x+0.1,y+0.06,3.0,0.36,size=11,bold=True,color=color)
    txt(s,text,x+0.1,y+0.5,6.1,0.44,size=9,color=LGRAY)
bar(s)

# ══════════════════════════════════════════════════════════════════
# SLIDE 15 — Summary
# ══════════════════════════════════════════════════════════════════
s=blank(); bg(s); num(s,15)
rect(s,0,0,13.33,0.68,BLUE)
txt(s,"PRISMA Review Agent — Architecture Summary",0.3,0.1,12.7,0.52,size=22,bold=True,color=NAVY,align=PP_ALIGN.CENTER)

pillars=[
    ("15-Step Pipeline",BLUE,
     "PRISMAReviewPipeline orchestrates everything. Calls tools, manages state, tracks PRISMA flow counts, assembles result."),
    ("9 Specialised LLM Agents",GREEN,
     "One agent per PRISMA step. All use pydantic-ai typed outputs + defer_model_check=True for any OpenRouter model."),
    ("Two-Stage Screening",YELLOW,
     "Agent 2 runs twice: INCLUSIVE T/A screening (batch 15) then STRICT FT screening (batch 10). PRISMA best practice."),
    ("Domain-Agnostic Design",ORANGE,
     "All clinical knowledge in ReviewProtocol. 11 RoB tools, any PICO, any outcome — zero code changes between reviews."),
    ("GRADE + RoB",PURPLE,
     "Agent 3 assesses RoB per study (11 tools). Agent 6 produces GRADE certainty per outcome. Both in final report."),
    ("Parallel Synthesis",TEAL,
     "asyncio.gather(): narrative + bias summary + GRADE x3 + limitations simultaneously. Report complete in one pass."),
]

for i,(title,color,text) in enumerate(pillars):
    col=i%3; row=i//3
    x=0.2+col*4.35; y=0.85+row*2.35
    rect(s,x,y,4.2,2.22,DGRAY,color,Pt(2))
    rect(s,x,y,4.2,0.5,color)
    txt(s,title,x+0.1,y+0.07,4.0,0.38,size=13,bold=True,color=NAVY)
    txt(s,text,x+0.1,y+0.62,4.0,1.48,size=11.5,color=LGRAY)

rect(s,0.2,5.55,12.95,0.8,DGRAY,BLUE,Pt(1.2))
txt(s,"ReviewProtocol  ->  Agent 1 (search queries)  ->  PubMed + bioRxiv Tools  ->  Dedup  ->  "
    "Agent 2 T/A Screening (inclusive)  ->  PMC Full-text  ->  Agent 2 FT Screening (strict)  ->  "
    "Agents 9/4/3 per article  ->  asyncio.gather Agents 5/7/6/8  ->  PRISMAReviewResult  ->  PRISMA 2020 Report",
    0.35,5.62,12.65,0.68,size=10.5,color=WHITE,align=PP_ALIGN.CENTER)

txt(s,"PRISMA Review Agent v1.0.0   |   Tek Raj Chhetri   |   Apache 2.0   |   2026   |   pip install synthscholar",
    0.3,6.6,12.7,0.38,size=12,color=LBLUE,align=PP_ALIGN.CENTER)
bar(s)

OUT="/Users/tekrajchhetri/Documents/research_codes_papers_writing/synthscholar/PRISMA_Agent_Architecture.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
